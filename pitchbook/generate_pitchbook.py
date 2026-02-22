"""Generate Newport Wholesalers Government Contract Entry Pitchbook v3.

Creates a 14-slide professional PowerPoint presentation targeting
Newport leadership to secure buy-in for government contract entry,
with Still Mind Creative as the operational partner.

v3 — Grounded in reality: no inflated TAMs, no salesy dream-selling,
no EBITDA multiples, no pipe-dream valuations. Each claim backed by real data.
Sent alongside a working dashboard with example contracts.

Usage:
    python pitchbook/generate_pitchbook.py
"""

import math
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = PROJECT_ROOT / "pitchbook"

# --- Color Palette ---
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
GOLD = RGBColor(0xD4, 0xA8, 0x43)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x7A, 0x33)
LIGHT_NAVY = RGBColor(0x2D, 0x3E, 0x6E)
GOLD_LIGHT = RGBColor(0xFF, 0xF8, 0xE7)
AMBER = RGBColor(0xE6, 0x8A, 0x00)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ========== HELPERS ==========

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text="",
                font_size=18, font_color=DARK_GRAY, bold=False,
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=16, font_color=DARK_GRAY, bold=False,
                          alignment=PP_ALIGN.LEFT, line_spacing=Pt(6)):
    """Add a text box with multiple styled paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = alignment
        p.space_after = line_spacing
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, font_color=DARK_GRAY, bullet_char="\u2022",
                    spacing=Pt(8)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Calibri"
        p.space_after = spacing

    return txBox


def add_kpi_card(slide, left, top, width, height, value, label,
                 value_color=NAVY, value_size=36, border_color=GOLD):
    """Add a KPI card with large value and small label."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    shape.shadow.inherit = False

    txBox = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.15),
        width - Inches(0.2), Inches(0.7)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(value_size)
    p.font.color.rgb = value_color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(
        left + Inches(0.1), top + height - Inches(0.55),
        width - Inches(0.2), Inches(0.45)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = MED_GRAY
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    return shape


def add_table(slide, left, top, width, rows, cols, data,
              col_widths=None, header_bg=NAVY, header_fg=WHITE,
              font_size=12):
    """Add a styled table."""
    row_height = Inches(0.4)
    table_height = row_height * rows
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, table_height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"

                if r == 0:
                    paragraph.font.color.rgb = header_fg
                    paragraph.font.bold = True
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def add_divider(slide, left, top, width, color=GOLD):
    """Add a horizontal gold divider line."""
    shape = slide.shapes.add_shape(1, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_callout_bar(slide, left, top, width, height, text,
                    bg_color=NAVY, text_color=GOLD, font_size=16):
    """Add a colored bar with centered text."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    add_textbox(slide, left, top, width, height,
                text, font_size=font_size, font_color=text_color,
                bold=True, alignment=PP_ALIGN.CENTER)
    return shape


def add_bar_chart(slide, left, top, width, height, categories, values,
                  series_name="", bar_color=NAVY, number_format=None):
    """Add a column chart with styled bars and optional data labels."""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = bar_color
    if number_format:
        series.has_data_labels = True
        labels = series.data_labels
        labels.font.size = Pt(12)
        labels.font.bold = True
        labels.font.color.rgb = NAVY
        labels.font.name = "Calibri"
        labels.number_format = number_format
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(11)
    cat_axis.tick_labels.font.name = "Calibri"
    cat_axis.tick_labels.font.color.rgb = DARK_GRAY
    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(10)
    val_axis.tick_labels.font.name = "Calibri"
    val_axis.tick_labels.font.color.rgb = MED_GRAY
    val_axis.has_major_gridlines = True
    val_axis.major_gridlines.format.line.color.rgb = LIGHT_GRAY
    return chart_frame


# ========== SLIDE BUILDERS ==========

def slide_01_cover(prs):
    """Slide 1: Cover — clean, professional, no inflated promises."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
                "Newport Wholesalers", font_size=48, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
                "A Strategy for Government Food Procurement",
                font_size=28, font_color=GOLD,
                alignment=PP_ALIGN.CENTER)

    add_divider(slide, Inches(4), Inches(4.3), Inches(5))

    add_textbox(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.5),
                "Prepared by Still Mind Creative LLC",
                font_size=18, font_color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(5.4), Inches(10), Inches(0.4),
                datetime.now().strftime("%B %Y"),
                font_size=16, font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(6.6), Inches(10), Inches(0.4),
                "CONFIDENTIAL", font_size=12, font_color=MED_GRAY,
                alignment=PP_ALIGN.CENTER)


def slide_02_you_built_something_real(prs):
    """Slide 2: 'You Built Something Real' — emotional foundation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "You Built Something Real", font_size=36, font_color=NAVY, bold=True)
    add_divider(slide, Inches(0.8), Inches(1.0), Inches(3))

    # LEFT: Vertical timeline — visual mass of 30 years
    # Stack of navy bars representing decades of continuous operation
    for i in range(15):
        y_bar = Inches(1.4) + i * Inches(0.22)
        bar_w = Inches(3.2) - i * Inches(0.03)
        alpha = max(0.4, 1.0 - i * 0.04)
        shape = slide.shapes.add_shape(1, Inches(0.8), y_bar, bar_w, Inches(0.14))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY if i < 10 else LIGHT_NAVY
        shape.line.fill.background()

    add_textbox(slide, Inches(0.8), Inches(4.8), Inches(3.5), Inches(0.6),
                "30 YEARS", font_size=42, font_color=NAVY, bold=True)
    add_textbox(slide, Inches(0.8), Inches(5.4), Inches(3.5), Inches(0.4),
                "of continuous operation", font_size=16, font_color=MED_GRAY)

    # RIGHT: Three bold statements with gold left borders
    statements = [
        "Real warehouses. Real trucks.\nReal employees.",
        "~30 years of auditable\nfinancial records.",
        "An American business that pays\nAmerican taxes.",
    ]
    for i, stmt in enumerate(statements):
        y_stmt = Inches(1.5) + i * Inches(1.4)
        # Gold left accent bar
        border = slide.shapes.add_shape(
            1, Inches(5.2), y_stmt, Inches(0.08), Inches(1.0))
        border.fill.solid()
        border.fill.fore_color.rgb = GOLD
        border.line.fill.background()
        add_textbox(slide, Inches(5.5), y_stmt + Inches(0.1),
                    Inches(7), Inches(0.8),
                    stmt, font_size=22, font_color=NAVY, bold=True)

    # Bottom callout — full width
    add_callout_bar(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8),
                    "In the current environment, this is the most valuable thing "
                    "a food distributor can have.",
                    font_size=16)


def slide_03_why_government_customers(prs):
    """Slide 3: 'Why Government Customers Are Different' — structural advantages."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Why Government Customers Are Different",
                font_size=36, font_color=NAVY, bold=True)
    add_divider(slide, Inches(0.8), Inches(1.0), Inches(3))

    # Four KPI cards showing structural advantages
    cards = [
        ("Net-30\nPayment", "Enforced by federal law\n(Prompt Payment Act)"),
        ("3-5 Year\nContracts", "With built-in\nprice escalators"),
        ("Recession-\nProof", "Spending INCREASED\nduring 2008-09"),
        ("Higher\nMargins", "5-15% above equivalent\nprivate sector"),
    ]
    card_w = Inches(2.7)
    card_h = Inches(1.4)
    gap = Inches(0.3)
    start_x = Inches(0.8)
    for i, (val, label) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_kpi_card(slide, x, Inches(1.3), card_w, card_h,
                     val, label, value_color=NAVY, value_size=22)

    # Bar chart: Federal food spending FY2023-2025 (USASpending data)
    add_textbox(slide, Inches(0.8), Inches(2.9), Inches(6), Inches(0.4),
                "Federal Food Spending Growth (USASpending.gov)",
                font_size=14, font_color=NAVY, bold=True)
    add_bar_chart(slide,
                  Inches(0.8), Inches(3.3), Inches(5.5), Inches(2.8),
                  categories=["FY2023", "FY2024", "FY2025"],
                  values=(1.276, 1.322, 1.467),
                  series_name="Federal Food Spending ($B)",
                  bar_color=NAVY,
                  number_format='"$"0.000"B"')

    # Right side: Context text
    context_lines = [
        "This is federal spending only on food",
        "wholesale NAICS codes (4244xx).",
        "",
        "State and local spending (school",
        "districts, county jails, state prisons)",
        "adds multiples more.",
        "",
        "Government is the structurally ideal",
        "client. They sign long contracts,",
        "print money, and pay on time.",
    ]
    add_multiline_textbox(slide, Inches(7.0), Inches(3.3),
                          Inches(5.5), Inches(2.8),
                          context_lines, font_size=13,
                          font_color=MED_GRAY, line_spacing=Pt(2))

    # Bottom note
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                "Source: USASpending.gov, NAICS 424410/424420/424430/424440/424450/424460/424470/424480/424490 + 722310. "
                "Federal food procurement only. State/local adds $9-13B+ annually.",
                font_size=10, font_color=MED_GRAY)


def slide_04_fraud_crisis(prs):
    """Slide 4: 'Why Now — The Fraud Crisis' — door is open for real companies."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Why Now \u2014 The Fraud Crisis",
                font_size=36, font_color=NAVY, bold=True)
    add_divider(slide, Inches(0.8), Inches(1.0), Inches(3))

    # Fraud data cards — sourced, not inflated
    fraud_items = [
        ("$6.8B", "False Claims Act recoveries\nFY2025 (all-time record)", "DOJ"),
        ("$250M+", "Feeding Our Future\nschool nutrition fraud stolen", "DOJ"),
        ("1,091", "8(a) firms suspended\n(25% of program)", "SBA OIG\nJan 2026"),
        ("Class-Action", "Aramark Dec 2025 \u2014\ncutting meals, contamination", "Reuters"),
        ("Voided", "Trinity OK contract\nvoided after 4 weeks", "OK DOC\nRecords"),
    ]
    card_w = Inches(2.3)
    card_h = Inches(2.0)
    gap = Inches(0.2)
    start_x = Inches(0.5)

    for i, (number, desc, source) in enumerate(fraud_items):
        x = start_x + i * (card_w + gap)
        y = Inches(1.3)

        # Card background
        card = slide.shapes.add_shape(1, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = RED
        card.line.width = Pt(1.5)

        # Number header bar
        header = slide.shapes.add_shape(1, x, y, card_w, Inches(0.5))
        header.fill.solid()
        header.fill.fore_color.rgb = NAVY
        header.line.fill.background()
        add_textbox(slide, x, y, card_w, Inches(0.5),
                    number, font_size=18, font_color=GOLD, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, x + Inches(0.1), y + Inches(0.6),
                    card_w - Inches(0.2), Inches(0.9),
                    desc, font_size=11, font_color=DARK_GRAY,
                    alignment=PP_ALIGN.CENTER)

        # Source
        add_textbox(slide, x + Inches(0.1), y + card_h - Inches(0.4),
                    card_w - Inches(0.2), Inches(0.3),
                    f"Source: {source}", font_size=9, font_color=MED_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Comparison table: "What Agencies Now Verify" vs "What Newport Has"
    add_textbox(slide, Inches(0.8), Inches(3.6), Inches(11), Inches(0.4),
                "What Agencies Now Verify vs. What Newport Has",
                font_size=16, font_color=NAVY, bold=True)

    comp_data = [
        ["What Agencies Now Verify", "What Newport Has"],
        ["Continuous operating history", "~30 years of continuous Florida operations"],
        ["Physical infrastructure", "Warehouses, trucks, cold chain (inspectable)"],
        ["Clean financial records", "Decades of auditable books and tax returns"],
        ["American ownership", "American-owned, American-operated, American taxes"],
    ]
    add_table(slide, Inches(1.5), Inches(4.1), Inches(10.3),
              len(comp_data), 2, comp_data,
              col_widths=[Inches(4.5), Inches(5.8)],
              font_size=13)

    # Bottom callout
    add_callout_bar(slide, Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.8),
                    "The door is open for real companies. Not because the market "
                    "is big \u2014 because agencies are desperate for legitimate, "
                    "verifiable vendors.",
                    font_size=15)


def slide_05_competitive_moat(prs):
    """Slide 5: 'Newport's Competitive Moat' — diagram + excluded competitors."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Newport\u2019s Competitive Moat",
                font_size=36, font_color=NAVY, bold=True)
    add_divider(slide, Inches(0.8), Inches(1.0), Inches(3))

    # --- MOAT DIAGRAM ---
    # Background moat rectangle (light gold)
    moat_x, moat_y = Inches(0.5), Inches(1.3)
    moat_w, moat_h = Inches(12.3), Inches(3.2)
    moat_bg = slide.shapes.add_shape(5, moat_x, moat_y, moat_w, moat_h)  # ROUNDED_RECT
    moat_bg.fill.solid()
    moat_bg.fill.fore_color.rgb = GOLD_LIGHT
    moat_bg.line.color.rgb = GOLD
    moat_bg.line.width = Pt(2)

    # Central box: NEWPORT
    cw, ch = Inches(3.2), Inches(0.8)
    cx = moat_x + moat_w // 2 - cw // 2
    cy = moat_y + moat_h // 2 - ch // 2
    center = slide.shapes.add_shape(5, cx, cy, cw, ch)
    center.fill.solid()
    center.fill.fore_color.rgb = NAVY
    center.line.fill.background()
    add_textbox(slide, cx, cy, cw, ch,
                "NEWPORT\nWHOLESALERS", font_size=16, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Surrounding advantage items (3 top, 3 bottom)
    items_top = [
        "~30 Years\nOperating History",
        "Real Infrastructure\n(Inspectable)",
        "Small Business\nQualification",
    ]
    items_bottom = [
        "Florida Home Base\n(Local Preference)",
        "Broad-Line\nProduct Range",
        "Existing Supply Chain\n& Relationships",
    ]

    item_w = Inches(3.2)
    item_h = Inches(0.7)
    item_gap = (moat_w - 3 * item_w - Inches(0.6)) // 2

    for i, text in enumerate(items_top):
        ix = moat_x + Inches(0.3) + i * (item_w + item_gap)
        iy = moat_y + Inches(0.25)
        box = slide.shapes.add_shape(5, ix, iy, item_w, item_h)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = GOLD
        box.line.width = Pt(1.5)
        add_textbox(slide, ix, iy, item_w, item_h,
                    text, font_size=12, font_color=NAVY, bold=True,
                    alignment=PP_ALIGN.CENTER)

    for i, text in enumerate(items_bottom):
        ix = moat_x + Inches(0.3) + i * (item_w + item_gap)
        iy = moat_y + moat_h - item_h - Inches(0.25)
        box = slide.shapes.add_shape(5, ix, iy, item_w, item_h)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = GOLD
        box.line.width = Pt(1.5)
        add_textbox(slide, ix, iy, item_w, item_h,
                    text, font_size=12, font_color=NAVY, bold=True,
                    alignment=PP_ALIGN.CENTER)

    # --- EXCLUDED SECTION ---
    add_textbox(slide, Inches(0.8), Inches(4.7), Inches(11.5), Inches(0.4),
                "EXCLUDED from small business set-asides:",
                font_size=16, font_color=RED, bold=True)

    excluded = [
        ("Sysco", "$76B"),
        ("US Foods", "$37B"),
        ("Aramark", "$17B"),
        ("Compass Group", "$40B"),
        ("Sodexo", "$25B"),
    ]
    ex_x = Inches(0.8)
    for name, rev in excluded:
        ex_w = Inches(2.0)
        shape = slide.shapes.add_shape(1, ex_x, Inches(5.15), ex_w, Inches(0.45))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = RED
        shape.line.width = Pt(1)
        add_textbox(slide, ex_x, Inches(5.15), ex_w, Inches(0.45),
                    f"\u2717 {name} ({rev})", font_size=11, font_color=RED,
                    bold=True, alignment=PP_ALIGN.CENTER)
        ex_x += ex_w + Inches(0.3)

    # Set-aside callout
    callout = slide.shapes.add_shape(
        1, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.3))
    callout.fill.solid()
    callout.fill.fore_color.rgb = GOLD_LIGHT
    callout.line.color.rgb = GOLD
    callout.line.width = Pt(2)

    add_textbox(slide, Inches(1.0), Inches(5.95), Inches(11.3), Inches(0.4),
                "~23% of federal contracts are reserved for small businesses. "
                "Newport qualifies.",
                font_size=16, font_color=NAVY, bold=True)
    add_textbox(slide, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.7),
                "The three largest food distributors in America cannot bid on these "
                "contracts. They are legally excluded. This is not a nice-to-have \u2014 "
                "it is a structural advantage that cannot be replicated by any "
                "billion-dollar competitor.",
                font_size=13, font_color=DARK_GRAY)


def slide_06_targets_identified(prs):
    """Slide 6: 'What We're Targeting' — realistic TAM first, then pursuit process."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "What We\u2019re Targeting \u2014 Realistic and Defensible",
                font_size=34, font_color=NAVY, bold=True)
    add_divider(slide, Inches(0.8), Inches(1.0), Inches(3))

    # TOP: Realistic TAM — contracts we can actually win
    add_textbox(slide, Inches(0.5), Inches(1.15), Inches(12), Inches(0.35),
                "Contracts Newport Can Fulfill \u2014 Direct Government and Private Sector Serving Government (FL + Nationwide)",
                font_size=14, font_color=NAVY, bold=True)

    tam_data = [
        ["Contract Type", "Available Volume", "Value Range", "Year 1 Target"],
        ["FL School District Category ITBs\n(meat, dairy, frozen, produce)",
         "67 districts (top 15 by\nenrollment targeted)",
         "$2-10M per\ncategory bid",
         "1-2 category wins"],
        ["BOP Quarterly Subsistence\n(6 FL federal prisons)",
         "Quarterly cycle \u2014\nconstant bid flow",
         "$75K-$2.5M\nper quarter",
         "1-2 quarterly wins"],
        ["FL County Jail Food Supply\n(self-operated kitchens)",
         "5-10 self-op counties\n(Pinellas, Polk, Brevard...)",
         "$1-3.5M/year\nper county",
         "1 win"],
        ["Private Contractor Supply\n(GEO, Trinity, CoreCivic)",
         "FL-based HQs \u2014\napply as vendor",
         "Facility-\ndependent",
         "Vendor approval\n+ first orders"],
        ["Nationwide SAM.gov Simplified\nAcquisitions ($15K-$350K)",
         "1,600-3,000 food postings\nper year nationwide",
         "$15K-$350K\neach",
         "5-10 responses"],
        ["USDA AMS / LFPA / FEMA\n(commodity, emergency, food bank)",
         "Multiple programs\nyear-round",
         "$200K-$5M\neach",
         "Applications\nsubmitted"],
    ]
    add_table(slide, Inches(0.2), Inches(1.55), Inches(12.9),
              len(tam_data), 4, tam_data,
              col_widths=[Inches(3.2), Inches(3.2), Inches(2.0), Inches(2.5)],
              font_size=10)

    # BOTTOM: How We Pursue Them — the process
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12), Inches(0.35),
                "How We Go After Them:",
                font_size=14, font_color=NAVY, bold=True)

    process_data = [
        ["Step", "What Happens", "Status"],
        ["1. Register", "SAM.gov, MFMP, VendorLink, DemandStar \u2014 get in the system", "Ready to start"],
        ["2. Monitor", "Daily scanning of all portals + saved search alerts", "Automated dashboard built"],
        ["3. Source Contacts", "Compile decision maker info from public directories, agency sites", "76 FL contacts identified ($0)"],
        ["4. Respond to Sources Sought", "Shape future solicitations \u2014 get known before bids drop", "Process defined"],
        ["5. Prepare Compliance", "Capability statement, certifications, financial docs, insurance", "Checklist ready"],
        ["6. Submit Bids", "RFQs, ITBs, RFPs \u2014 start with micro-purchases and simplified", "Templates in progress"],
        ["7. Track & Iterate", "Measure results, assess often, test, learn, adjust", "Dashboard included"],
    ]
    add_table(slide, Inches(0.2), Inches(5.1), Inches(12.9),
              len(process_data), 3, process_data,
              col_widths=[Inches(2.5), Inches(6.7), Inches(3.7)],
              font_size=10)

    add_textbox(slide, Inches(0.8), Inches(7.1), Inches(11.5), Inches(0.3),
                "Same process scales to all 50 states. The dashboard sent alongside "
                "this presentation demonstrates the scanning and scoring system.",
                font_size=11, font_color=NAVY, bold=True,
                alignment=PP_ALIGN.CENTER)


def slide_07_how_door_opens(prs):
    """Slide 7: 'How the Door Opens' — three graduated threshold panels."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "How the Door Opens",
                font_size=36, font_color=NAVY, bold=True)
    add_divider(slide, Inches(0.8), Inches(1.0), Inches(3))

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
                "Government has specific thresholds that create low-barrier entry "
                "points. Start with the smallest door.",
                font_size=14, font_color=MED_GRAY)

    # Three door panels — increasing size left to right, top-aligned
    doors = [
        {
            "num": "1",
            "title": "Micro-Purchase",
            "threshold": "Under $15,000",
            "height": Inches(3.2),
            "lines": [
                "No competition required.",
                "Agency buys directly.",
                "Just show up and deliver.",
                "",
                "This is how you get your",
                "first past performance",
                "reference.",
            ],
        },
        {
            "num": "2",
            "title": "Simplified Acquisition",
            "threshold": "$15K \u2013 $350K",
            "height": Inches(3.8),
            "lines": [
                "Only 3 quotes needed.",
                "Average 2-4 bidders.",
                "Most vendors don't bother.",
                "",
                "Posted on SAM.gov daily.",
                "1,600-3,000 food postings",
                "per year nationwide.",
                "",
                "Newport's sweet spot.",
            ],
        },
        {
            "num": "3",
            "title": "Full Competition",
            "threshold": "$350K+",
            "height": Inches(4.5),
            "lines": [
                "Formal proposals required.",
                "Past performance evaluated.",
                "Price + technical capability.",
                "",
                "You EARN this by stacking",
                "wins from Doors 1 and 2.",
                "",
                "This is where the multi-year,",
                "multi-million-dollar contracts",
                "live. School districts, BOP",
                "quarterly subsistence, county",
                "jail supply.",
            ],
        },
    ]

    widths = [Inches(3.2), Inches(3.5), Inches(4.0)]
    door_y = Inches(1.6)
    x = Inches(0.8)

    for i, door in enumerate(doors):
        w = widths[i]
        h = door["height"]

        # Door frame (rounded rectangle)
        shape = slide.shapes.add_shape(5, x, door_y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = GOLD
        shape.line.width = Pt(3)

        # Door number badge (circle)
        badge_sz = Inches(0.45)
        badge = slide.shapes.add_shape(
            9, x + w // 2 - badge_sz // 2,  # OVAL
            door_y + Inches(0.15), badge_sz, badge_sz)
        badge.fill.solid()
        badge.fill.fore_color.rgb = NAVY
        badge.line.fill.background()
        add_textbox(slide, x + w // 2 - badge_sz // 2,
                    door_y + Inches(0.15), badge_sz, badge_sz,
                    door["num"], font_size=18, font_color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Title
        add_textbox(slide, x + Inches(0.1), door_y + Inches(0.7),
                    w - Inches(0.2), Inches(0.35),
                    door["title"], font_size=16, font_color=NAVY, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Threshold
        add_textbox(slide, x + Inches(0.1), door_y + Inches(1.05),
                    w - Inches(0.2), Inches(0.3),
                    door["threshold"], font_size=13, font_color=GOLD, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Body lines
        y_line = door_y + Inches(1.4)
        for line in door["lines"]:
            add_textbox(slide, x + Inches(0.15), y_line,
                        w - Inches(0.3), Inches(0.25),
                        line, font_size=11, font_color=DARK_GRAY,
                        alignment=PP_ALIGN.CENTER)
            y_line += Inches(0.22)

        # Arrow between doors
        if i < 2:
            arrow_x = x + w + Inches(0.05)
            add_textbox(slide, arrow_x, door_y + Inches(1.2),
                        Inches(0.4), Inches(0.4),
                        "\u27A4", font_size=24, font_color=GOLD, bold=True,
                        alignment=PP_ALIGN.CENTER)

        x += w + Inches(0.45)

    # Bottom note
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
                "The visual progression IS the strategy: start left, move right "
                "over time. Each threshold crossed builds the credentials for "
                "the next level.",
                font_size=12, font_color=NAVY, bold=True,
                alignment=PP_ALIGN.CENTER)


def slide_08_realistic_path(prs):
    """Slide 8: 'The Realistic Path — Year 1 Through Year 5' — staircase chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "The Realistic Path \u2014 Year 1 Through Year 5",
                font_size=34, font_color=NAVY, bold=True)
    add_divider(slide, Inches(0.8), Inches(1.0), Inches(3))

    # Staircase chart using shaped rectangles (full control over labels)
    years_data = [
        ("Year 1", "$500K-$2M", 1.25, NAVY,
         "Stack 15-25 micro +\nsimplified wins. FL focus\n+ nationwide SAM.gov.\nBuild credentials."),
        ("Year 2", "$2-5M", 3.5, NAVY,
         "Past performance from\nYear 1 unlocks $100K-$500K\ncontracts. First BOP wins.\nFirst school district."),
        ("Year 3", "$5-10M", 7.5, LIGHT_NAVY,
         "Full competition bids.\nMulti-facility contracts.\nSE regional expansion.\nFSMC vendor channel."),
        ("Year 4", "$8-15M", 11.5, LIGHT_NAVY,
         "Established vendor status.\nOption year renewals.\nLarger contract vehicles."),
        ("Year 5", "$10-20M", 15.0, GOLD,
         "Government division is a\ncore business line. Repeat\ncustomers. Sticky,\nrecurring revenue."),
    ]

    max_val = 15.0
    max_bar_h = Inches(2.5)
    bar_w = Inches(2.0)
    gap = Inches(0.3)
    start_x = Inches(0.8)
    base_y = Inches(4.1)  # Bottom of bars

    for i, (year, range_label, val, color, desc) in enumerate(years_data):
        x = start_x + i * (bar_w + gap)
        bar_h = int(max_bar_h * val / max_val)
        if bar_h < Inches(0.4):
            bar_h = Inches(0.4)
        y = base_y - bar_h

        # Bar
        bar = slide.shapes.add_shape(1, x, y, bar_w, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Range label on bar
        fc = WHITE if color != GOLD else NAVY
        add_textbox(slide, x, y + Inches(0.05), bar_w, Inches(0.35),
                    range_label, font_size=14, font_color=fc, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Year label below bar
        add_textbox(slide, x, base_y + Inches(0.05), bar_w, Inches(0.3),
                    year, font_size=13, font_color=NAVY, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Description below year label
        add_textbox(slide, x, base_y + Inches(0.35), bar_w, Inches(1.2),
                    desc, font_size=10, font_color=DARK_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Connecting note below chart
    add_textbox(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.35),
                "Each step unlocks the next. Past performance from small wins is "
                "the credential that opens larger doors.",
                font_size=13, font_color=NAVY, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Acquisition mention — one sentence, passing, at the bottom
    add_textbox(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.8),
                "As government contracts become a recurring revenue stream, they "
                "also strengthen Newport's position for any future strategic "
                "decisions \u2014 including potential acquisition interest, if "
                "that matters down the road.",
                font_size=11, font_color=MED_GRAY,
                alignment=PP_ALIGN.CENTER)


def slide_09_year_one_machine(prs):
    """Slide 9: 'Year 1: Build the Machine' — three phases + success metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Year 1: Build the Machine",
                font_size=36, font_color=NAVY, bold=True)
    add_divider(slide, Inches(0.8), Inches(1.0), Inches(3))

    # Three phase blocks — horizontal timeline
    phases = [
        {
            "title": "MONTHS 1-2",
            "subtitle": "Foundation",
            "color": NAVY,
            "items": [
                "SAM.gov registration",
                "MFMP registration",
                "Capability statement",
                "Saved searches live",
                "Dashboard built",
            ],
        },
        {
            "title": "MONTHS 3-6",
            "subtitle": "First Wins",
            "color": LIGHT_NAVY,
            "items": [
                "Respond to 10+ solicitations",
                "Sources Sought responses",
                "First 3-5 micro/simplified wins",
                "First past performance refs",
            ],
        },
        {
            "title": "MONTHS 7-12",
            "subtitle": "Scale",
            "color": GOLD,
            "items": [
                "15+ total contracts",
                "USDA AMS application",
                "FSMC vendor applications",
                "Nationwide scanning active",
                "Proposal library built",
            ],
        },
    ]

    phase_w = Inches(3.7)
    phase_gap = Inches(0.3)
    phase_x = Inches(0.5)
    phase_y = Inches(1.2)
    phase_h = Inches(3.0)

    for i, phase in enumerate(phases):
        x = phase_x + i * (phase_w + phase_gap)

        # Phase header bar
        hdr = slide.shapes.add_shape(1, x, phase_y, phase_w, Inches(0.5))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = phase["color"]
        hdr.line.fill.background()
        hdr_fc = WHITE if phase["color"] != GOLD else NAVY
        add_textbox(slide, x, phase_y, phase_w, Inches(0.5),
                    f'{phase["title"]}  |  {phase["subtitle"]}',
                    font_size=14, font_color=hdr_fc, bold=True,
                    alignment=PP_ALIGN.CENTER)

        # Phase content area
        body = slide.shapes.add_shape(
            1, x, phase_y + Inches(0.5), phase_w, phase_h - Inches(0.5))
        body.fill.solid()
        body.fill.fore_color.rgb = LIGHT_GRAY
        body.line.color.rgb = phase["color"]
        body.line.width = Pt(1.5)

        # Bullet items
        add_bullet_list(slide, x + Inches(0.15),
                        phase_y + Inches(0.6),
                        phase_w - Inches(0.3), phase_h - Inches(0.7),
                        phase["items"], font_size=12, font_color=DARK_GRAY,
                        bullet_char="\u2713", spacing=Pt(6))

        # Arrow between phases
        if i < 2:
            arrow_x = x + phase_w + Inches(0.02)
            add_textbox(slide, arrow_x, phase_y + Inches(1.0),
                        Inches(0.3), Inches(0.4),
                        "\u27A4", font_size=20, font_color=GOLD, bold=True,
                        alignment=PP_ALIGN.CENTER)

    # Bottom: Year 1 SUCCESS METRICS (KPI cards, in priority order)
    add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.35),
                "YEAR 1 SUCCESS METRICS (in priority order):",
                font_size=14, font_color=NAVY, bold=True)

    metrics = [
        ("5+", "Past Performance\nReferences"),
        ("15+", "Solicitations\nResponded To"),
        ("10+", "Agency\nRelationships"),
        ("$500K-$2M", "Revenue\n(secondary)"),
    ]
    met_w = Inches(2.7)
    met_gap = Inches(0.25)
    met_x = Inches(0.8)

    for i, (val, label) in enumerate(metrics):
        x = met_x + i * (met_w + met_gap)
        border = GOLD if i == 0 else MED_GRAY
        vc = GREEN if i == 0 else NAVY
        add_kpi_card(slide, x, Inches(4.9), met_w, Inches(1.2),
                     val, label, value_color=vc, value_size=26,
                     border_color=border)

    # Highlight first metric
    add_textbox(slide, met_x, Inches(6.2), met_w, Inches(0.3),
                "\u2B06 THIS is the real asset",
                font_size=10, font_color=GREEN, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, met_x + 3 * (met_w + met_gap), Inches(6.2),
                met_w, Inches(0.3),
                "listed LAST \u2014 deliberately",
                font_size=10, font_color=MED_GRAY,
                alignment=PP_ALIGN.CENTER)


def slide_10_small_wins_compound(prs):
    """Slide 10: 'The Power of Compound Interest' — leverage deals, contract stacking."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "The Power of Compound Interest",
                font_size=36, font_color=NAVY, bold=True)
    add_divider(slide, Inches(0.8), Inches(1.0), Inches(3))

    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.4),
                "Government contracts don\u2019t just add revenue \u2014 they stack. "
                "Each win is a credential that unlocks the next. This is compound "
                "interest applied to business development.",
                font_size=14, font_color=DARK_GRAY)

    # How contracts stack over time
    stack_data = [
        ["", "Year 1", "Year 2", "Year 3", "Year 5"],
        ["New contracts won", "3-5", "5-8", "4-12", "5-10"],
        ["Carried from prior years\n(multi-year + renewals)", "\u2014", "3-5", "8-10", "20+"],
        ["TOTAL ACTIVE CONTRACTS", "3-5", "8-13", "12-22", "25-30+"],
        ["Estimated annual revenue", "$500K-$2M", "$2-5M", "$5-10M", "$10-20M"],
        ["Avg contract length", "Quarterly-1yr", "1-3 years", "3-5 years", "3-5 years"],
        ["Incumbent retention rate", "N/A (new)", "70-80%", "75-85%", "80-90%"],
    ]
    add_table(slide, Inches(0.5), Inches(1.7), Inches(12.3),
              len(stack_data), 5, stack_data,
              col_widths=[Inches(3.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)],
              font_size=11)

    # The leverage mechanism
    add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.35),
                "Why This Works \u2014 The Leverage Mechanism:",
                font_size=14, font_color=NAVY, bold=True)

    leverage_data = [
        ["What You Leverage", "Into What", "Why It Compounds"],
        ["1 flawless micro-purchase delivery",
         "Documented past performance reference",
         "Agencies weight past performance 20-30%\nin evaluations \u2014 you now have what\nmost new entrants don\u2019t"],
        ["3-5 past performance references",
         "Eligibility for $100K-$500K contracts",
         "You\u2019ve bought your government track\nrecord \u2014 the most valuable credential\nin procurement"],
        ["Multi-quarter delivery track record",
         "Incumbent advantage on renewals",
         "Agencies don\u2019t switch vendors without\ncause \u2014 retention runs 70-90%"],
        ["~30 years of general business history",
         "Credibility no new entrant can match",
         "You\u2019re not building from zero \u2014 you\u2019re\ntranslating an existing track record\ninto a new market"],
    ]
    add_table(slide, Inches(0.3), Inches(5.0), Inches(12.7),
              len(leverage_data), 3, leverage_data,
              col_widths=[Inches(3.2), Inches(3.5), Inches(4.0)],
              font_size=10)

    # Bottom callout
    add_callout_bar(slide, Inches(1.0), Inches(6.8), Inches(11.3), Inches(0.6),
                    "Newport has ~30 years of general track record. Now we buy the "
                    "government track record \u2014 and we\u2019re well-positioned to do "
                    "it because of those ~30 years.",
                    font_size=14)


def slide_11_division_of_labor(prs):
    """Slide 11: 'A Suggested Approach' — here's what I'd recommend + how I can help."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "A Suggested Approach",
                font_size=36, font_color=NAVY, bold=True)
    add_divider(slide, Inches(0.8), Inches(1.0), Inches(3))

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.4),
                "Based on everything above, here\u2019s what I\u2019d recommend "
                "\u2014 and if you\u2019re open to it, here\u2019s how I can support you.",
                font_size=14, font_color=DARK_GRAY)

    # LEFT: What I'd Recommend Newport Do
    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
                "What I\u2019d Recommend Newport Do",
                font_size=18, font_color=NAVY, bold=True)
    recommend_items = [
        "Register on SAM.gov, MFMP, VendorLink, DemandStar",
        "Assign 1 existing salesperson (8-10 hrs/week)",
        "Attend 1-2 industry conferences for face-to-face relationships",
        "Start with micro-purchases and simplified acquisitions",
        "Review and approve all proposals before submission",
        "Deliver flawlessly on every contract won",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5),
                    recommend_items, font_size=13, font_color=DARK_GRAY,
                    bullet_char="\u2713", spacing=Pt(8))

    # RIGHT: How I Can Support You
    add_textbox(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.4),
                "How I Can Support You",
                font_size=18, font_color=NAVY, bold=True)
    support_items = [
        "Daily monitoring of all procurement portals",
        "Opportunity scoring and prioritization",
        "Proposal drafting and response preparation",
        "Compliance tracking and deadline management",
        "Past performance documentation after every delivery",
        "Weekly pipeline and progress reports",
        "Continuous research and strategy adjustment",
    ]
    add_bullet_list(slide, Inches(6.8), Inches(2.1), Inches(5.5), Inches(2.8),
                    support_items, font_size=13, font_color=DARK_GRAY,
                    bullet_char="\u2713", spacing=Pt(7))

    # What Newport Does NOT Need — still useful context
    add_textbox(slide, Inches(0.8), Inches(4.7), Inches(11), Inches(0.35),
                "What Newport Does NOT Need to Learn or Hire For:",
                font_size=14, font_color=RED, bold=True)

    no_items = [
        "Procurement rules",
        "Portal monitoring",
        "Proposal writing",
        "Compliance tracking",
        "Specialist hires",
    ]
    x_row_start = Inches(0.8)
    item_w = Inches(2.3)
    for i, item in enumerate(no_items):
        x = x_row_start + i * item_w
        shape = slide.shapes.add_shape(1, x, Inches(5.1), item_w - Inches(0.1),
                                       Inches(0.45))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.fill.background()
        add_textbox(slide, x, Inches(5.1), item_w - Inches(0.1), Inches(0.45),
                    f"\u2717 {item}", font_size=12, font_color=MED_GRAY,
                    bold=True, alignment=PP_ALIGN.CENTER)

    # Bottom: The real ask — no retainers, just cover costs
    callout = slide.shapes.add_shape(
        1, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2))
    callout.fill.solid()
    callout.fill.fore_color.rgb = GOLD_LIGHT
    callout.line.color.rgb = GOLD
    callout.line.width = Pt(2)

    add_textbox(slide, Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.4),
                "No retainers. No consulting fees.",
                font_size=18, font_color=NAVY, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.6),
                "Just cover the direct costs of trying it \u2014 registrations, "
                "materials, and tools \u2014 and let me bring value to your business. "
                "If it works, we both win.",
                font_size=14, font_color=DARK_GRAY,
                alignment=PP_ALIGN.CENTER)


def slide_12_what_it_costs(prs):
    """Slide 12: 'What It Actually Costs' — comprehensive, honest cost breakdown."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "What It Actually Costs \u2014 Everything",
                font_size=36, font_color=NAVY, bold=True)
    add_divider(slide, Inches(0.8), Inches(1.0), Inches(3))

    # LEFT: Fixed / Startup Costs
    add_textbox(slide, Inches(0.3), Inches(1.15), Inches(6.5), Inches(0.35),
                "Implementation Costs (Year 1)",
                font_size=14, font_color=NAVY, bold=True)

    fixed_data = [
        ["Item", "Cost"],
        ["SAM.gov registration (UEI number)", "Free"],
        ["MFMP + VendorLink + DemandStar", "$0-$400/yr"],
        ["Industry memberships (ACFSA, FSNA)", "$225/yr"],
        ["Conference attendance (1 person, 1-2 events)", "$1,500-$3,000"],
        ["Capability statement + sales materials", "$300-$500"],
        ["SQF/HACCP certification (if not held)", "$3,000-$8,000"],
        ["Tech tools (scanning, automation, AI)", "$2,400-$6,000/yr"],
        ["Insurance endorsements (gov contract rider)", "$500-$2,000/yr"],
        ["Legal review (contract templates, first bids)", "$1,000-$3,000"],
        ["Compliance documentation (gov-specific SOPs)", "$500-$1,500"],
        ["TOTAL YEAR 1 FIXED", "$10,000-$25,000"],
    ]
    add_table(slide, Inches(0.2), Inches(1.5), Inches(6.4),
              len(fixed_data), 2, fixed_data,
              col_widths=[Inches(4.2), Inches(2.2)],
              font_size=10)

    # RIGHT: Variable / Per-Contract Costs
    add_textbox(slide, Inches(6.8), Inches(1.15), Inches(5.8), Inches(0.35),
                "Per-Contract Costs (As You Win)",
                font_size=14, font_color=NAVY, bold=True)

    variable_data = [
        ["Item", "Cost"],
        ["Bid/performance bonds (>$35K contracts)", "1-3% of contract value"],
        ["Background checks (BOP/military access)", "$50-$100/person"],
        ["Proposal prep (printing, binding, shipping)", "$100-$300/bid"],
        ["Delivery compliance (temp logs, documentation)", "Operational cost"],
    ]
    add_table(slide, Inches(6.8), Inches(1.5), Inches(6.2),
              len(variable_data), 2, variable_data,
              col_widths=[Inches(3.8), Inches(2.4)],
              font_size=10)

    # What one win returns
    add_textbox(slide, Inches(6.8), Inches(3.5), Inches(5.8), Inches(0.35),
                "What One Win Returns",
                font_size=14, font_color=GREEN, bold=True)

    return_data = [
        ["Contract Type", "Value"],
        ["Micro-purchase contract", "$5K-$15K"],
        ["Simplified acquisition", "$15K-$350K"],
        ["School district category bid", "$100K-$350K/year"],
        ["BOP quarterly subsistence", "$75K-$2.5M/quarter"],
        ["Gross margin (all types)", "15-20%"],
    ]
    add_table(slide, Inches(6.8), Inches(3.9), Inches(6.2),
              len(return_data), 2, return_data,
              col_widths=[Inches(3.8), Inches(2.4)],
              font_size=10, header_bg=GREEN)

    # ROI ratio — positioned below tables, no overlap
    ratio_y = Inches(5.7)
    cost_box = slide.shapes.add_shape(
        5, Inches(1.0), ratio_y, Inches(3.5), Inches(0.8))
    cost_box.fill.solid()
    cost_box.fill.fore_color.rgb = LIGHT_GRAY
    cost_box.line.color.rgb = NAVY
    cost_box.line.width = Pt(2)
    add_textbox(slide, Inches(1.0), ratio_y + Inches(0.05),
                Inches(3.5), Inches(0.3),
                "Year 1 All-In Cost", font_size=11, font_color=MED_GRAY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), ratio_y + Inches(0.3),
                Inches(3.5), Inches(0.4),
                "$10K-$25K", font_size=22, font_color=NAVY, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(4.7), ratio_y + Inches(0.1),
                Inches(0.5), Inches(0.6),
                "vs", font_size=14, font_color=MED_GRAY, bold=True,
                alignment=PP_ALIGN.CENTER)

    ret_box = slide.shapes.add_shape(
        5, Inches(5.4), ratio_y, Inches(6.5), Inches(0.8))
    ret_box.fill.solid()
    ret_box.fill.fore_color.rgb = GOLD_LIGHT
    ret_box.line.color.rgb = GREEN
    ret_box.line.width = Pt(2)
    add_textbox(slide, Inches(5.4), ratio_y + Inches(0.05),
                Inches(6.5), Inches(0.3),
                "One Simplified Acquisition Win", font_size=11, font_color=MED_GRAY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.4), ratio_y + Inches(0.3),
                Inches(6.5), Inches(0.4),
                "$50K-$200K", font_size=22, font_color=GREEN, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Bottom note
    add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                "No retainers. No consulting fees. Cover the costs of trying it. "
                "Payback: 1 contract.",
                font_size=14, font_color=NAVY, bold=True,
                alignment=PP_ALIGN.CENTER)


def slide_13_what_could_go_wrong(prs):
    """Slide 13: 'What Could Go Wrong' — genuine risk analysis, agile mindset."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "What Could Go Wrong \u2014 And How We Mitigate It",
                font_size=32, font_color=NAVY, bold=True)
    add_divider(slide, Inches(0.8), Inches(1.0), Inches(3))

    # Risk table — 6 real business risks, centered
    risk_data = [
        ["Risk", "What Happens", "How Bad?", "Mitigation"],
        [
            "Take on a contract\nwe can't deliver",
            "Negative past performance\nrating in CPARS. Damages\nfuture bids for years.",
            "SERIOUS\n#1 risk \u2014 one bad\nrating follows\nyou 3+ years",
            "Start with micro-purchases and small\ncontracts ONLY. Never bid on anything\nNewport can't deliver with current\ninfrastructure. Scale up only after proven.",
        ],
        [
            "Learning curve \u2014\ntakes time to win\nfirst contracts",
            "~$10-25K invested, salesperson\ntime allocated, revenue\ntakes time to materialize.",
            "MANAGEABLE\nNewport doesn't\nneed this to\nsurvive",
            "We research, assess often, test, learn,\nadjust. Everything is measured, methodical,\nintentional. 6-month checkpoints. It's not\nif \u2014 it's how long until we find the way.",
        ],
        [
            "Underprice a\ncontract and lose\nmoney on delivery",
            "Negative margin on a\ncontract Newport must\nfulfill.",
            "MODERATE\nBounded to the\ncontract size.\n$50K at -5% = $2.5K",
            "Build pricing templates with full cost\nmodeling. Never bid without complete cost\nanalysis. Small contracts first = cheap\nlessons before the stakes get larger.",
        ],
        [
            "Government\npayment delays\ndespite Net-30",
            "Cash flow strain if large\ncontract payments are late.",
            "LOW\nFederal agencies\nare generally\nreliable",
            "Start with federal contracts (strongest\npayment enforcement). Don't take contract\nsizes that would strain cash flow if\npayment runs 60 days instead of 30.",
        ],
        [
            "Compliance failure\n(food safety, docs)",
            "Contract termination,\npotential debarment.",
            "SEVERE if it\nhappens \u2014 but\nNewport already\nhas infrastructure",
            "Ensure HACCP/SQF certs current. Create\ncompliance checklist for every delivery.\nSame discipline Newport already applies\nfor private sector with government wrapper.",
        ],
        [
            "Opportunity cost \u2014\nsalesperson time\ndiverted",
            "8-10 hrs/week not spent\non existing customers.",
            "REAL but\nBOUNDED\nRedirect back\nafter 6-12mo",
            "Use existing salesperson capacity, not a\nnew hire. All admin work is handled. 6-month\ncheckpoint evaluates if time is well spent.\nThis is additive, not a bet-the-farm move.",
        ],
    ]
    add_table(slide, Inches(0.5), Inches(1.2), Inches(12.3),
              len(risk_data), 4, risk_data,
              col_widths=[Inches(2.1), Inches(2.6), Inches(2.1), Inches(3.8)],
              font_size=9)

    # Bottom callout
    add_callout_bar(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.2),
                    "The biggest risk is taking on a contract you can't deliver. "
                    "The mitigation is simple: start small, deliver perfectly, "
                    "scale only when proven.\n"
                    "Every other risk is bounded and manageable. Newport doesn\u2019t "
                    "need government contracts to survive \u2014 which is itself a "
                    "major advantage.\nWe can be patient, methodical, and selective "
                    "about which opportunities to pursue.",
                    font_size=13)


def slide_14_what_happens_next(prs):
    """Slide 14: 'What I'd Suggest as Next Steps' — soft close, value-first."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Inches(1.5), Inches(0.2), Inches(10), Inches(0.6),
                "What I\u2019d Suggest as Next Steps",
                font_size=36, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_divider(slide, Inches(4), Inches(0.85), Inches(5), color=GOLD)

    # Section 1: "If You're Open to It"
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
                "If You\u2019re Open to Exploring This, Here\u2019s What I\u2019d Build:",
                font_size=18, font_color=GOLD, bold=True)

    build_items = [
        "Complete opportunity tracking system (scanning SAM.gov, MFMP, "
        "DemandStar \u2014 dashboard demo included with this presentation)",
        "Capability statement and compliance package",
        "Pricing templates for micro-purchase and simplified acquisition bids",
        "Pipeline with scoring, deadlines, and status tracking",
        "Regular monitoring and progress reports \u2014 what\u2019s working, "
        "what we adjust",
    ]
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(build_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2713  {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_after = Pt(4)

    # Section 2: Year 1 Targets — left side
    sec2_y = Inches(3.3)
    add_textbox(slide, Inches(0.8), sec2_y, Inches(5.5), Inches(0.35),
                "Year 1 Target:", font_size=18, font_color=GOLD, bold=True)

    targets = [
        "15-25 contracts won (micro + simplified)",
        "$500K-$2M revenue",
        "5+ documented past performance references",
        "Active pipeline in FL + nationwide",
    ]
    txBox2 = slide.shapes.add_textbox(
        Inches(0.8), sec2_y + Inches(0.35), Inches(5.5), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(targets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_after = Pt(4)

    # Section 3: Realistic Cost — right side
    add_textbox(slide, Inches(6.8), sec2_y, Inches(5.5), Inches(0.35),
                "Realistic Year 1 Cost:", font_size=18, font_color=GOLD, bold=True)

    costs = [
        "$10,000-$25,000 (registrations, materials, insurance,",
        "    legal, tools, conferences)",
        "1 salesperson, 8-10 hours/week (existing staff)",
        "6-month checkpoints to assess and adjust",
    ]
    txBox3 = slide.shapes.add_textbox(
        Inches(6.8), sec2_y + Inches(0.35), Inches(5.5), Inches(1.3))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    for i, item in enumerate(costs):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = f"\u2022  {item}" if not item.startswith("    ") else f"   {item.strip()}"
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_after = Pt(3)

    # Divider before closing
    add_divider(slide, Inches(2), Inches(5.3), Inches(9), color=GOLD)

    # Closing — soft, genuine
    add_textbox(slide, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.5),
                "No retainers. No fees. Just cover the costs of trying it and "
                "let me bring value to your business.",
                font_size=16, font_color=GOLD, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.6),
                "The dashboard is ready. The research is done. "
                "I\u2019m ready when you are.",
                font_size=22, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
                "Still Mind Creative LLC  |  " + datetime.now().strftime("%B %Y")
                + "  |  Confidential",
                font_size=11, font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ========== MAIN ==========

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Building pitchbook v3 slides...")
    slides = [
        ("Slide 1:  Cover", slide_01_cover),
        ("Slide 2:  You Built Something Real", slide_02_you_built_something_real),
        ("Slide 3:  Why Government Customers Are Different", slide_03_why_government_customers),
        ("Slide 4:  Why Now \u2014 The Fraud Crisis", slide_04_fraud_crisis),
        ("Slide 5:  Newport's Competitive Moat", slide_05_competitive_moat),
        ("Slide 6:  The Targets Are Already Identified", slide_06_targets_identified),
        ("Slide 7:  How the Door Opens", slide_07_how_door_opens),
        ("Slide 8:  The Realistic Path (Year 1-5)", slide_08_realistic_path),
        ("Slide 9:  Year 1: Build the Machine", slide_09_year_one_machine),
        ("Slide 10: How Small Wins Compound", slide_10_small_wins_compound),
        ("Slide 11: You Sell. We Handle the Bureaucracy.", slide_11_division_of_labor),
        ("Slide 12: What It Actually Costs", slide_12_what_it_costs),
        ("Slide 13: What Could Go Wrong", slide_13_what_could_go_wrong),
        ("Slide 14: What Happens Next (Close)", slide_14_what_happens_next),
    ]

    for name, builder in slides:
        print(f"  {name}")
        builder(prs)

    output_path = OUTPUT_DIR / f"newport_pitchbook_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    prs.save(str(output_path))
    print(f"\nPitchbook saved to: {output_path}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
